import os
import logging

import docx
import fitz
import nltk
from docx.opc.exceptions import OpcError, PackageNotFoundError as DocxPackageNotFoundError
from nltk.tokenize import sent_tokenize, word_tokenize
from pptx import Presentation
from pptx.exc import (
    InvalidXmlError,
    PackageNotFoundError as PptxPackageNotFoundError,
    PythonPptxError,
)

_NLTK_SETUP_DONE = False
logger = logging.getLogger(__name__)


class FileParsingError(ValueError):
    """Raised when a local file cannot be parsed into article data."""


def _setup_nltk():
    global _NLTK_SETUP_DONE
    if _NLTK_SETUP_DONE:
        return

    for resource, pkg in [
        ("tokenizers/punkt_tab", "punkt_tab"),
        ("tokenizers/punkt",     "punkt"),
    ]:
        try:
            nltk.data.find(resource)
        except (LookupError, OSError):
            nltk.download(pkg, quiet=True)

    _NLTK_SETUP_DONE = True


def _ensure_local_file(file_path: str) -> str:
    normalized_path = os.path.expanduser(file_path)

    if not os.path.exists(normalized_path):
        raise FileParsingError(f"Input file not found: {file_path}")
    if not os.path.isfile(normalized_path):
        raise FileParsingError(f"Input path is not a file: {file_path}")

    return normalized_path


def _read_txt_file(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError as exc:
        raise FileParsingError(
            f"Could not decode text file as UTF-8: {file_path}"
        ) from exc
    except OSError as exc:
        raise FileParsingError(f"Could not read text file '{file_path}': {exc}") from exc


def _read_docx_file(file_path: str) -> str:
    try:
        document = docx.Document(file_path)
    except DocxPackageNotFoundError as exc:
        raise FileParsingError(f"DOCX file not found or unreadable: {file_path}") from exc
    except OpcError as exc:
        raise FileParsingError(f"Invalid or corrupt DOCX file: {file_path}") from exc
    except OSError as exc:
        raise FileParsingError(f"Could not read DOCX file '{file_path}': {exc}") from exc

    return " ".join(p.text for p in document.paragraphs)


def _read_pptx_file(file_path: str) -> str:
    try:
        presentation = Presentation(file_path)
    except PptxPackageNotFoundError as exc:
        raise FileParsingError(f"PPTX file not found or unreadable: {file_path}") from exc
    except (PythonPptxError, InvalidXmlError) as exc:
        raise FileParsingError(f"Invalid or corrupt PPTX file: {file_path}") from exc
    except OSError as exc:
        raise FileParsingError(f"Could not read PPTX file '{file_path}': {exc}") from exc

    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            if not text_frame:
                continue
            for paragraph in text_frame.paragraphs:
                if paragraph.text:
                    texts.append(paragraph.text)

    return " ".join(texts)


def _read_pdf_file(file_path: str) -> str:
    try:
        doc = fitz.open(file_path)
    except fitz.FileNotFoundError as exc:
        raise FileParsingError(f"PDF file not found: {file_path}") from exc
    except (fitz.EmptyFileError, fitz.FileDataError, RuntimeError, ValueError) as exc:
        raise FileParsingError(f"Invalid or corrupt PDF file: {file_path}") from exc
    except OSError as exc:
        raise FileParsingError(f"Could not read PDF file '{file_path}': {exc}") from exc

    try:
        return " ".join(page.get_text() for page in doc)
    except RuntimeError as exc:
        raise FileParsingError(f"Failed to extract text from PDF: {file_path}") from exc
    finally:
        doc.close()


def parse_local_file(file_path: str) -> dict:
    normalized_path = _ensure_local_file(file_path)
    ext = os.path.splitext(normalized_path)[1].lower()

    if ext == ".txt":
        raw_text = _read_txt_file(normalized_path)
    elif ext == ".docx":
        raw_text = _read_docx_file(normalized_path)
    elif ext == ".pptx":
        raw_text = _read_pptx_file(normalized_path)
    elif ext == ".pdf":
        raw_text = _read_pdf_file(normalized_path)
    else:
        raise FileParsingError(
            f"Unsupported file format: {ext}. Supported: .txt, .docx, .pptx, .pdf"
        )

    raw_text = " ".join((raw_text or "").split())
    try:
        _setup_nltk()
        sentences = sent_tokenize(raw_text)
    except (LookupError, OSError, ValueError) as exc:
        raise FileParsingError(
            f"Failed to tokenize content extracted from '{file_path}': {exc}"
        ) from exc

    data = []
    for sentence in sentences:
        if not sentence.strip():
            continue
        try:
            words = word_tokenize(sentence)
        except (LookupError, OSError, ValueError) as exc:
            raise FileParsingError(
                f"Failed to tokenize a sentence extracted from '{file_path}': {exc}"
            ) from exc
        data.append({"sentence": sentence, "words": words})

    article_data = {
        "title": os.path.splitext(os.path.basename(normalized_path))[0],
        "data": data,
    }

    logger.info(
        f"Parsed local file '{normalized_path}' into {len(article_data['data'])} sentences."
    )
    return article_data
